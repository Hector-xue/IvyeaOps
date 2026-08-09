"""交付物导出。

这份文件的存在理由是**可被追问**：每条结论后面跟着指标、数值、时间窗和来源。
所以测试盯的不是"能不能生成文件"，而是"证据有没有真的落到纸上"。
"""
from __future__ import annotations

import pytest

from app.core.findings import Action, Evidence, Finding, FindingList
from app.services import deliverable

openpyxl = pytest.importorskip("openpyxl")


def _sample() -> dict:
    fl = FindingList(
        findings=[
            Finding(
                title="B0ABC 的 SP 广告 30 天花了 820 美元零出单",
                severity="critical",
                priority_score=95,
                reasoning="点击量足够大，转化为零，不是样本不足的问题。",
                evidence=[
                    Evidence(metric="spend_30d", value=820.5, unit="USD",
                             target="B0ABC", source="ads_api", as_of="2026-07-01~07-30"),
                    Evidence(metric="orders_30d", value=0, target="B0ABC",
                             source="ads_api", as_of="2026-07-01~07-30"),
                ],
                actions=[Action(type="pause_campaign", target="SP-B0ABC",
                                detail="暂停该广告活动", guardrail="点击≥15 且 0 单",
                                reversible=True, confidence=0.9)],
            ),
            Finding(
                title="建议整体下调竞价",
                severity="low",
                actions=[Action(type="adjust_bid", detail="全局降 10%",
                                reversible=False, confidence=0.3)],
            ),
        ],
        data_notes="仅覆盖 SP，未含 SB/SD。",
    )
    payload = fl.model_dump()
    payload["unsupported"] = fl.unsupported()
    return payload


def test_evidence_lands_on_its_own_sheet(tmp_path):
    """证据独立成页，是这份东西和"AI 生成的一段话"最大的区别。"""
    out = deliverable.build_xlsx(tmp_path / "d.xlsx", _sample(), {"job_id": "J1"})
    wb = openpyxl.load_workbook(out)
    assert wb.sheetnames == ["结论", "证据", "说明"]

    ev = wb["证据"]
    flat = [[c.value for c in row] for row in ev.iter_rows(min_row=2)]
    metrics = {r[2] for r in flat}
    assert "spend_30d" in metrics and "orders_30d" in metrics
    spend = next(r for r in flat if r[2] == "spend_30d")
    assert spend[3] == 820.5              # 数值本身，不是"很高"
    assert spend[6] == "ads_api"          # 来源
    assert "2026-07-01~07-30" in str(spend[7])   # 时间窗


def test_a_finding_without_evidence_is_marked_not_silently_blank(tmp_path):
    """空证据必须显式写出来。留白会让人以为只是"碰巧没填"，
    而这恰恰是最该被质疑的一条。"""
    out = deliverable.build_xlsx(tmp_path / "d.xlsx", _sample(), {})
    ev = openpyxl.load_workbook(out)["证据"]
    marks = [c.value for row in ev.iter_rows(min_row=2) for c in row]
    assert "（无结构化证据）" in marks


def test_most_severe_finding_comes_first(tmp_path):
    out = deliverable.build_xlsx(tmp_path / "d.xlsx", _sample(), {})
    ws = openpyxl.load_workbook(out)["结论"]
    assert ws.cell(row=2, column=2).value == "紧急"
    assert "820" in str(ws.cell(row=2, column=4).value)


def test_irreversible_action_is_flagged(tmp_path):
    """用户决定"照做还是不照做"时，最先要知道的是做错了能不能撤回来。"""
    out = deliverable.build_xlsx(tmp_path / "d.xlsx", _sample(), {})
    ws = openpyxl.load_workbook(out)["结论"]
    col = {ws.cell(row=1, column=c).value: c for c in range(1, 9)}
    flags = [ws.cell(row=r, column=col["可回滚"]).value for r in (2, 3)]
    assert set(flags) == {"是", "否"}


def test_empty_findings_says_so_instead_of_producing_a_blank_sheet(tmp_path):
    """没结论时给一张空表，用户会以为是导出坏了。"""
    out = deliverable.build_xlsx(tmp_path / "d.xlsx",
                                 {"findings": [], "data_notes": ""}, {})
    ws = openpyxl.load_workbook(out)["结论"]
    assert "没有得出可支撑的结论" in str(ws.cell(row=2, column=4).value)


def test_info_sheet_records_where_this_came_from(tmp_path):
    out = deliverable.build_xlsx(tmp_path / "d.xlsx", _sample(),
                                 {"job_id": "J7", "kind": "广告审计"})
    info = openpyxl.load_workbook(out)["说明"]
    blob = " ".join(str(c.value) for row in info.iter_rows() for c in row)
    assert "J7" in blob and "广告审计" in blob
    assert "仅覆盖 SP" in blob                      # 数据边界要跟着交付物走
    assert "建议整体下调竞价" in blob                # 无证据的结论要点名


def test_markdown_carries_the_same_evidence():
    md = deliverable.build_markdown(_sample(), {"job_id": "J1", "kind": "广告审计"})
    assert "spend_30d = 820.5USD" in md
    assert "来源：ads_api" in md
    assert "护栏：点击≥15 且 0 单" in md
    assert "可回滚：否" in md                        # 第二条动作
    assert "这条没有结构化证据支撑" in md


def test_markdown_labels_target_as_object_not_goal():
    """target 是「这条证据说的是哪个对象」，不是对比目标值 ——
    标成"目标 B0ABC"会让读者以为那是个要达成的数。"""
    md = deliverable.build_markdown(_sample(), {})
    assert "对象 B0ABC" in md
    assert "目标 B0ABC" not in md


# ── PPTX：拿去开会的那一份 ──────────────────────────────────────────────
def test_pptx_puts_evidence_on_the_same_slide_as_the_conclusion(tmp_path):
    """会上被问到"凭什么"的时候，翻页找附录就已经晚了。"""
    pptx = pytest.importorskip("pptx")
    out = deliverable.build_pptx(tmp_path / "d.pptx", _sample(), {"kind": "广告审计"})
    prs = pptx.Presentation(str(out))
    # 封面 + 2 条结论 + 数据说明
    assert len(prs.slides) == 4
    second = " ".join(sh.text_frame.text for sh in prs.slides[1].shapes if sh.has_text_frame)
    assert "820" in second and "spend_30d" in second      # 证据就在这一页
    assert "护栏" in second and "点击≥15" in second


def test_pptx_most_severe_first(tmp_path):
    pytest.importorskip("pptx")
    import pptx as _p
    out = deliverable.build_pptx(tmp_path / "d.pptx", _sample(), {})
    prs = _p.Presentation(str(out))
    second = " ".join(sh.text_frame.text for sh in prs.slides[1].shapes if sh.has_text_frame)
    assert "紧急" in second


def test_pptx_never_pins_a_latin_font(tmp_path):
    """一旦写死某个西文字体，中文会在放映的机器上退化成方块。
    让 PowerPoint/WPS/Keynote 用自己的默认字体，中文一定显示得出来。"""
    pytest.importorskip("pptx")
    import pptx as _p
    out = deliverable.build_pptx(tmp_path / "d.pptx", _sample(), {})
    prs = _p.Presentation(str(out))
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    assert run.font.name is None, f"钉死了字体：{run.font.name}"


def test_pptx_carries_the_data_boundary(tmp_path):
    pytest.importorskip("pptx")
    import pptx as _p
    out = deliverable.build_pptx(tmp_path / "d.pptx", _sample(), {})
    prs = _p.Presentation(str(out))
    last = " ".join(sh.text_frame.text for sh in prs.slides[-1].shapes if sh.has_text_frame)
    assert "仅覆盖 SP" in last
    assert "建议整体下调竞价" in last     # 无证据的结论要点名


# ── PDF：有 chrome 才做，没有要说清楚 ───────────────────────────────────
def test_pdf_without_chrome_raises_a_translatable_error(tmp_path, monkeypatch):
    """没装 chrome **不是错误**，是这台机器上没有。要能翻译成"你可以怎么办"，
    而不是变成一个 500。"""
    monkeypatch.setattr(deliverable, "chrome_bin", lambda: "")
    with pytest.raises(RuntimeError, match="没有找到"):
        deliverable.build_pdf(tmp_path / "x.pdf", "<html>hi</html>")


def test_pdf_is_produced_when_chrome_can_actually_run(tmp_path):
    """**"装了 Chrome"不等于"能跑无头"**。macOS 的 CI runner 上就是这样：
    二进制在，但无头打印会卡死。所以这里把"跑不动"当成环境不具备而跳过，
    而不是当成回归 —— 真正的回归是"能跑却没产出 PDF"，下面那两条断言管这个。
    """
    if not deliverable.chrome_bin():
        pytest.skip("这台机器没有 chrome")
    try:
        out = deliverable.build_pdf(
            tmp_path / "x.pdf",
            "<html><meta charset='utf-8'><body><h1>广告优化方案</h1>"
            "<p>花费 820.50 USD</p></body></html>")
    except RuntimeError as exc:
        pytest.skip(f"这个环境跑不了无头 Chrome：{exc}")
    blob = out.read_bytes()
    assert blob.startswith(b"%PDF")
    assert len(blob) > 1000
