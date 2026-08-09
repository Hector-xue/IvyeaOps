"""交付物导出：把结论变成能直接发给老板/客户的文件。

分析跑完了、结论也对，但用户下一步要做的事往往是"把这个发给别人"。到这一步，
一个网页是不够的 —— 他需要一份能落到邮件附件里、能打印、能在会上翻的东西。

为什么以 Excel 为主
------------------
亚马逊运营的工作流本来就跑在表格上（后台导出的是表格、给老板的周报是表格、
和供应链对的是表格）。给一份 PDF 好看但改不动，给一份表格他能直接往里加自己的列。

**证据单独成页**是这份交付物和"AI 生成的一段话"最大的区别：每条结论后面跟着
指标、数值、时间窗和数据来源。别人质疑"你凭什么这么说"的时候，翻到第二页就行。
"""
from __future__ import annotations

import logging
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List

logger = logging.getLogger("ivyea.services.deliverable")

# 配色跟广告审计报表保持一致，同一个产品出来的东西该像一家的。
_HEADER_BG = "1F3864"
_EVIDENCE_BG = "EAF1FB"
_WARN_BG = "FCE4D6"

# core/findings 的 severity 词表（critical/high/medium/low）。
_SEVERITY_CN = {"critical": "紧急", "high": "高", "medium": "中", "low": "低"}
_SEVERITY_BG = {"critical": "F4B3B3", "high": "FCE4D6", "medium": "", "low": ""}


def _confidence_cn(action: Dict[str, Any]) -> str:
    """置信度在 **Action** 上，是 0~1 的小数（不是 Finding 上的字符串）。"""
    if not action:
        return ""
    try:
        v = float(action.get("confidence") or 0)
    except (TypeError, ValueError):
        return ""
    if v <= 0:
        return ""
    return f"{v * 100:.0f}%"


def _cell(ws, row: int, col: int, value: Any, *, bold=False, wrap=False,
          fill: str = "", color: str = ""):
    from openpyxl.styles import Alignment, Font, PatternFill
    c = ws.cell(row=row, column=col, value=value)
    if bold or color:
        c.font = Font(bold=bold, color=color or "000000")
    if wrap:
        c.alignment = Alignment(wrap_text=True, vertical="top")
    if fill:
        c.fill = PatternFill("solid", fgColor=fill)
    return c


def _widths(ws, widths: List[int]) -> None:
    from openpyxl.utils import get_column_letter
    for i, w in enumerate(widths, start=1):
        ws.column_dimensions[get_column_letter(i)].width = w


def build_xlsx(out_path: Path, findings: Dict[str, Any], meta: Dict[str, Any]) -> Path:
    """把一份 FindingList（见 core/findings）写成交付物。

    三页：结论 / 证据 / 说明。证据独立成页正是这份东西可被追问的地方。
    """
    from openpyxl import Workbook

    items: List[Dict[str, Any]] = list(findings.get("findings") or [])
    wb = Workbook()

    # ── 第 1 页：结论 ──────────────────────────────────────────────────
    ws = wb.active
    ws.title = "结论"
    _widths(ws, [6, 10, 8, 40, 34, 40, 24, 10])
    head = ["#", "严重度", "优先级", "结论", "推理", "建议动作", "护栏", "可回滚"]
    for i, h in enumerate(head, start=1):
        _cell(ws, 1, i, h, bold=True, fill=_HEADER_BG, color="FFFFFF")

    # 按严重度再按优先级排，让最该先看的落在第一行。
    order = {s: i for i, s in enumerate(("critical", "high", "medium", "low"))}
    items.sort(key=lambda f: (order.get(str(f.get("severity")), 9),
                              -int(f.get("priority_score") or 0)))

    row = 2
    for idx, f in enumerate(items, start=1):
        sev = str(f.get("severity") or "medium")
        actions = f.get("actions") or []
        first = actions[0] if actions else {}
        _cell(ws, row, 1, idx)
        _cell(ws, row, 2, _SEVERITY_CN.get(sev, sev), fill=_SEVERITY_BG.get(sev, ""))
        _cell(ws, row, 3, int(f.get("priority_score") or 0))
        _cell(ws, row, 4, f.get("title") or "", wrap=True)
        _cell(ws, row, 5, f.get("reasoning") or "", wrap=True)
        detail = first.get("detail") or first.get("type") or ""
        conf = _confidence_cn(first)
        _cell(ws, row, 6, f"{detail}（置信度 {conf}）" if conf else detail, wrap=True)
        _cell(ws, row, 7, first.get("guardrail") or "", wrap=True)
        # 「可回滚」给单独一列，是因为用户在决定"照做还是不照做"时，最先要
        # 知道的是做错了能不能撤回来。不可回滚的标黄，翻表时一眼能挑出来。
        rev = bool(first.get("reversible", True)) if first else True
        _cell(ws, row, 8, "是" if rev else "否", fill="" if rev else _WARN_BG)
        row += 1

    if not items:
        _cell(ws, 2, 4, "这次没有得出可支撑的结论。数据不足时不硬凑结论，"
                        "是这个工具刻意的行为，不是出错。", wrap=True)

    # ── 第 2 页：证据 ──────────────────────────────────────────────────
    ev = wb.create_sheet("证据")
    _widths(ev, [6, 40, 16, 14, 10, 14, 20, 30])
    head = ["结论#", "对应结论", "指标", "数值", "单位", "对象", "数据来源", "时间窗 / 备注"]
    for i, h in enumerate(head, start=1):
        _cell(ev, 1, i, h, bold=True, fill=_HEADER_BG, color="FFFFFF")

    r = 2
    for idx, f in enumerate(items, start=1):
        evidences = f.get("evidence") or []
        if not evidences:
            # **空证据必须显式写出来**，不能让这一页看起来只是"碰巧没填"。
            _cell(ev, r, 1, idx)
            _cell(ev, r, 2, f.get("title") or "", wrap=True)
            _cell(ev, r, 3, "（无结构化证据）", fill=_WARN_BG)
            r += 1
            continue
        for e in evidences:
            _cell(ev, r, 1, idx)
            _cell(ev, r, 2, f.get("title") or "", wrap=True)
            _cell(ev, r, 3, e.get("metric") or "")
            _cell(ev, r, 4, e.get("value"))
            _cell(ev, r, 5, e.get("unit") or "")
            _cell(ev, r, 6, e.get("target") or "")
            _cell(ev, r, 7, e.get("source") or "", fill=_EVIDENCE_BG)
            _cell(ev, r, 8, " / ".join(x for x in [e.get("as_of") or "", e.get("note") or ""] if x),
                  wrap=True)
            r += 1

    # ── 第 3 页：说明 ──────────────────────────────────────────────────
    info = wb.create_sheet("说明")
    _widths(info, [22, 70])
    rows = [
        ("生成时间", datetime.now().strftime("%Y-%m-%d %H:%M:%S")),
        ("来源任务", str(meta.get("job_id") or "")),
        ("分析类型", str(meta.get("kind") or "")),
        ("数据说明", str(findings.get("data_notes") or "")),
    ]
    unsupported = findings.get("unsupported") or []
    if unsupported:
        rows.append(("证据不足的结论",
                     "；".join(str(u) for u in unsupported)[:900]))
    rows += [
        ("怎么读这份东西",
         "「结论」页每条都对应「证据」页的若干行，证据里写明了指标、数值、"
         "时间窗和数据来源。有人质疑结论时，直接翻到证据页。"),
        ("免责",
         "结论由 AI 基于你提供的数据得出，建议动作请在执行前自行复核。"
         "标注为「不可回滚」的动作尤其如此。"),
        ("生成方",
         "IvyeaOps（自托管）。这份文件在你自己的机器上生成，数据没有离开过它。"),
    ]
    for i, (k, v) in enumerate(rows, start=1):
        _cell(info, i, 1, k, bold=True)
        _cell(info, i, 2, v, wrap=True)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(out_path))
    return out_path


# ── PPTX：给"要拿去开会"的那一步 ────────────────────────────────────────

def build_pptx(out_path: Path, findings: Dict[str, Any], meta: Dict[str, Any]) -> Path:
    """把结论做成一份能直接放的经营周报。

    每条结论一页，**证据写在同一页上**而不是另起附录 —— 会上被问到"凭什么"的
    时候，翻页找附录就已经晚了。

    字体刻意不指定：一旦写死某个西文字体，中文会在放映的机器上退化成方块。
    让 PowerPoint / WPS / Keynote 用它自己的默认字体，中文一定显示得出来。
    """
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    items = _sorted_items(findings)
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)   # 16:9

    # 封面
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _tb(slide, Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.2),
        str(meta.get("kind") or "分析结论"), size=40, bold=True)
    sub = f"{datetime.now().strftime('%Y-%m-%d')}"
    if meta.get("job_id"):
        sub += f" · 任务 {meta['job_id']}"
    _tb(slide, Inches(0.9), Inches(3.7), Inches(11.5), Inches(0.6), sub, size=16,
        color=RGBColor(0x66, 0x6E, 0x6A))
    _tb(slide, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.6),
        f"共 {len(items)} 条结论 · 由 IvyeaOps 在本机生成，数据未离开这台机器",
        size=12, color=RGBColor(0x8A, 0x92, 0x8E))

    for idx, f in enumerate(items, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sev = str(f.get("severity") or "medium")
        _tb(slide, Inches(0.7), Inches(0.5), Inches(1.6), Inches(0.45),
            _SEVERITY_CN.get(sev, sev), size=14, bold=True,
            color=RGBColor(0xC0, 0x39, 0x2B) if sev in ("critical", "high")
            else RGBColor(0x66, 0x6E, 0x6A))
        _tb(slide, Inches(0.7), Inches(1.0), Inches(11.9), Inches(1.1),
            f"{idx}. {f.get('title') or ''}", size=26, bold=True)
        if f.get("reasoning"):
            _tb(slide, Inches(0.7), Inches(2.1), Inches(11.9), Inches(0.9),
                str(f["reasoning"]), size=14, color=RGBColor(0x44, 0x4C, 0x48))

        lines = []
        for e in (f.get("evidence") or []):
            bits = f"{e.get('metric')} = {e.get('value')}{e.get('unit') or ''}"
            tail = " · ".join(x for x in [str(e.get('source') or ''), str(e.get('as_of') or '')] if x)
            lines.append(f"• {bits}" + (f"（{tail}）" if tail else ""))
        if not lines:
            lines = ["• 这条没有结构化证据支撑，仅供参考"]
        _tb(slide, Inches(0.7), Inches(3.1), Inches(7.4), Inches(3.4),
            "证据\n" + "\n".join(lines), size=13)

        acts = f.get("actions") or []
        first = acts[0] if acts else {}
        act_lines = [f"• {first.get('detail') or first.get('type') or '（无）'}"]
        if first.get("guardrail"):
            act_lines.append(f"• 护栏：{first['guardrail']}")
        act_lines.append(f"• 可回滚：{'是' if first.get('reversible', True) else '否'}")
        _tb(slide, Inches(8.4), Inches(3.1), Inches(4.2), Inches(3.4),
            "建议动作\n" + "\n".join(act_lines), size=13)

    # 数据边界单独一页：口径不跟着结论走，结论就没法被复核。
    if findings.get("data_notes") or findings.get("unsupported"):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _tb(slide, Inches(0.7), Inches(0.8), Inches(11.9), Inches(0.8),
            "数据说明", size=28, bold=True)
        body = str(findings.get("data_notes") or "")
        unsup = findings.get("unsupported") or []
        if unsup:
            body += ("\n\n以下结论没有结构化证据支撑：\n"
                     + "\n".join(f"• {u}" for u in unsup))
        _tb(slide, Inches(0.7), Inches(1.8), Inches(11.9), Inches(4.5), body, size=14)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


def _tb(slide, left, top, width, height, text: str, *, size=14, bold=False, color=None):
    from pptx.util import Pt
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(str(text).split("\n")):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        for run in para.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            if color is not None:
                run.font.color.rgb = color
    return box


def _sorted_items(findings: Dict[str, Any]) -> List[Dict[str, Any]]:
    items = list(findings.get("findings") or [])
    order = {s: i for i, s in enumerate(("critical", "high", "medium", "low"))}
    items.sort(key=lambda f: (order.get(str(f.get("severity")), 9),
                              -int(f.get("priority_score") or 0)))
    return items


# ── PDF：有 chrome 才做 ──────────────────────────────────────────────────

def chrome_bin() -> str:
    """找一个能用的 chrome。找不到返回空字符串。

    **刻意不把它写进依赖**：这是个要分发给普通卖家的自托管产品，为了一个导出
    格式就要求每台机器装 chrome，代价远大于收益。有就用，没有就告诉用户
    怎么自己拿到 PDF（HTML 报告在浏览器里 Ctrl+P 就是）。
    """
    import shutil
    for name in ("google-chrome", "chromium", "chromium-browser",
                 "google-chrome-stable", "msedge"):
        found = shutil.which(name)
        if found:
            return found
    for hard in (r"C:\Program Files\Google\Chrome\Application\chrome.exe",
                 "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"):
        if Path(hard).is_file():
            return hard
    return ""


def build_pdf(out_path: Path, html: str) -> Path:
    """用 headless chrome 把一份 HTML 打成 PDF。没有 chrome 时抛 RuntimeError，
    由调用方翻译成给用户看的说明。"""
    import subprocess
    import tempfile

    binary = chrome_bin()
    if not binary:
        raise RuntimeError("这台机器上没有找到 Chrome/Chromium")

    out_path.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory() as tmp:
        src = Path(tmp) / "report.html"
        src.write_text(html, encoding="utf-8")
        cmd = [binary, "--headless=new", "--disable-gpu", "--no-sandbox",
               f"--print-to-pdf={out_path}", "--no-pdf-header-footer",
               f"--user-data-dir={tmp}/profile", src.as_uri()]
        # **30 秒，不是 90 秒。** 这是一个用户点了按钮在等的动作 —— 让他干等
        # 一分半再看到失败，比早点告诉他"这台机器上不行、去 HTML 里打印"糟得多。
        # （macOS 的 CI runner 上实测过这条路会卡死，装了 Chrome 不等于能跑。）
        try:
            proc = subprocess.run(cmd, capture_output=True, timeout=30)
        except subprocess.TimeoutExpired as exc:
            raise RuntimeError(
                "Chrome 没能在 30 秒内生成 PDF（这台机器上的无头模式可能不可用）"
            ) from exc
    if not out_path.is_file() or out_path.stat().st_size == 0:
        tail = (proc.stderr or b"").decode("utf-8", "replace")[-300:]
        raise RuntimeError(f"Chrome 没有生成 PDF：{tail}")
    return out_path


def build_markdown(findings: Dict[str, Any], meta: Dict[str, Any]) -> str:
    """同一份内容的 Markdown 版本 —— 贴进飞书文档、Notion 或 issue 里用。"""
    items = list(findings.get("findings") or [])
    lines = [
        f"# 分析结论 · {meta.get('kind') or ''}".strip(),
        "",
        f"生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M')}"
        + (f" · 任务 {meta.get('job_id')}" if meta.get("job_id") else ""),
        "",
    ]
    if not items:
        lines += ["这次没有得出可支撑的结论。数据不足时不硬凑结论，是刻意的行为。", ""]

    for idx, f in enumerate(items, start=1):
        sev = str(f.get("severity") or "medium")
        lines.append(f"## {idx}. [{_SEVERITY_CN.get(sev, sev)}] {f.get('title') or ''}")
        if f.get("reasoning"):
            lines += ["", str(f["reasoning"])]
        ev = f.get("evidence") or []
        if ev:
            lines += ["", "**证据**", ""]
            for e in ev:
                bits = [f"{e.get('metric') or ''} = {e.get('value')}{e.get('unit') or ''}"]
                if e.get("target"):
                    # target 是「这条证据说的是哪个对象」（ASIN/关键词/广告活动），
                    # 不是对比目标值 —— 见 core/findings 的字段说明。
                    bits.append(f"对象 {e['target']}")
                if e.get("as_of"):
                    bits.append(str(e["as_of"]))
                if e.get("source"):
                    bits.append(f"来源：{e['source']}")
                lines.append("- " + " · ".join(bits))
        else:
            lines += ["", "> 这条没有结构化证据支撑，仅供参考。"]
        for a in (f.get("actions") or []):
            lines += ["", f"**建议动作**：{a.get('detail') or a.get('type') or ''}"]
            if a.get("guardrail"):
                lines.append(f"护栏：{a['guardrail']}")
            conf = _confidence_cn(a)
            lines.append(f"可回滚：{'是' if a.get('reversible', True) else '否'}"
                         + (f" · 置信度 {conf}" if conf else ""))
        lines.append("")

    if findings.get("data_notes"):
        lines += ["---", "", f"数据说明：{findings['data_notes']}", ""]
    lines.append("_由 IvyeaOps 在你自己的机器上生成，数据没有离开过它。_")
    return "\n".join(lines)
